#!/usr/bin/env python3
"""Презентация ЧОО «Гюрза» — об организации, фоны, переходы."""
from __future__ import annotations

import math
import random
from pathlib import Path

from PIL import Image, ImageDraw, ImageFilter
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml import parse_xml
from pptx.oxml.ns import nsdecls
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parent
SCREENS = ROOT / "assets" / "screens"
ASSETS = ROOT / "assets" / "generated"
OUT = ROOT / "gyurza-presentation.pptx"

W, H = 1920, 1080
BG = RGBColor(10, 14, 22)
PANEL = RGBColor(18, 24, 36)
GOLD = RGBColor(212, 175, 55)
GOLD_L = RGBColor(240, 217, 140)
TEXT = RGBColor(245, 247, 250)
MUTED = RGBColor(145, 155, 175)
CREAM = RGBColor(248, 246, 240)
INK = RGBColor(24, 28, 36)


def hex_vertex(cx: float, cy: float, r: float, deg: float) -> tuple[float, float]:
    rad = math.radians(deg)
    return cx + r * math.cos(rad), cy - r * math.sin(rad)


def draw_hex_fragment(
    draw: ImageDraw.ImageDraw,
    cx: float,
    cy: float,
    r: float,
    start_edge: int,
    edge_count: int,
    rotation: float,
    color: tuple[int, int, int],
    alpha: int,
    width: int = 3,
) -> None:
    """Кусочки соты — только часть граней, не целый шестиугольник."""
    pts = [hex_vertex(cx, cy, r, rotation + i * 60) for i in range(6)]
    for i in range(edge_count):
        idx = (start_edge + i) % 6
        nxt = (idx + 1) % 6
        draw.line([pts[idx], pts[nxt]], fill=(*color, alpha), width=width)


def draw_broken_hex_field(draw: ImageDraw.ImageDraw, w: int, h: int, seed: int) -> None:
    """Крупные фрагменты сот на чёрном фоне — как на сайте."""
    rng = random.Random(seed)
    gold = (212, 175, 55)

    placements = [
        (w * 0.82, h * 0.18, 220, 0, 2, 18, 2),
        (w * 0.78, h * 0.22, 180, 2, 3, 22, 3),
        (w * -0.05, h * 0.55, 260, 1, 2, 16, 2),
        (w * 0.12, h * 0.72, 200, 4, 3, 20, 2),
        (w * 0.92, h * 0.68, 240, 3, 2, 14, 2),
        (w * 0.55, h * 0.88, 190, 5, 2, 12, 2),
        (w * 0.35, h * -0.08, 170, 0, 3, 15, 2),
        (w * 0.68, h * 0.42, 150, 2, 2, 10, 1),
    ]
    for cx, cy, r, start, count, alpha, width in placements:
        draw_hex_fragment(draw, cx, cy, r, start, count, rng.uniform(-8, 8), gold, alpha, width)

    # Редкие полупрозрачные «островки» заливки
    for _ in range(4):
        cx = rng.uniform(w * 0.05, w * 0.95)
        cy = rng.uniform(h * 0.1, h * 0.9)
        r = rng.uniform(90, 160)
        rot = rng.uniform(0, 60)
        pts = [hex_vertex(cx, cy, r, rot + i * 60) for i in range(3)]
        draw.polygon(pts, fill=(212, 175, 55, rng.randint(4, 9)))


def make_backgrounds() -> dict[str, Path]:
    ASSETS.mkdir(parents=True, exist_ok=True)
    paths: dict[str, Path] = {}

    for name, seed in [("dark", 7), ("dark_alt", 42), ("dark_soft", 99)]:
        img = Image.new("RGBA", (W, H), (4, 5, 8, 255))
        draw = ImageDraw.Draw(img, "RGBA")
        draw_broken_hex_field(draw, W, H, seed)
        # Лёгкая вignette — без жёлтых полос
        vignette = Image.new("RGBA", (W, H), (0, 0, 0, 0))
        vd = ImageDraw.Draw(vignette)
        vd.rectangle((0, 0, W, H), fill=(0, 0, 0, 60))
        vd.ellipse((W // 2 - 900, H // 2 - 520, W // 2 + 900, H // 2 + 520), fill=(0, 0, 0, 0))
        img = Image.alpha_composite(img, vignette.filter(ImageFilter.GaussianBlur(90)))
        p = ASSETS / f"bg-{name}.png"
        img.convert("RGB").save(p, quality=95)
        paths[name.replace("dark_alt", "alt").replace("dark_soft", "soft").replace("dark", "dark")] = p

    paths = {"dark": ASSETS / "bg-dark.png", "alt": ASSETS / "bg-dark_alt.png", "soft": ASSETS / "bg-dark_soft.png"}
    return paths


def set_transition(slide, kind: str = "fade") -> None:
    el = slide._element
    for child in list(el):
        if child.tag.endswith("transition"):
            el.remove(child)
    if kind == "fade":
        xml = (
            f'<p:transition {nsdecls("p")} spd="med" advClick="1">'
            f'<p:fadeThroughBlack/></p:transition>'
        )
    elif kind == "push":
        xml = (
            f'<p:transition {nsdecls("p")} spd="med" advClick="1">'
            f'<p:push dir="l"/></p:transition>'
        )
    else:
        xml = f'<p:transition {nsdecls("p")} spd="med" advClick="1"><p:fade/></p:transition>'
    el.insert(2, parse_xml(xml))


def fill_bg(slide, path: Path) -> None:
    slide.shapes.add_picture(str(path), 0, 0, width=prs.slide_width, height=prs.slide_height)
    # Send to back
    pic = slide.shapes[-1]
    slide.shapes._spTree.remove(pic._element)
    slide.shapes._spTree.insert(2, pic._element)


def textbox(slide, left, top, width, height, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    box.text_frame.word_wrap = True
    box.text_frame.vertical_anchor = MSO_ANCHOR.TOP
    for p in box.text_frame.paragraphs:
        p.alignment = align
    return box


def add_text(slide, left, top, width, height, lines: list[tuple[str, int, RGBColor, bool]], align=PP_ALIGN.LEFT):
    box = textbox(slide, left, top, width, height, align)
    tf = box.text_frame
    for i, (txt, size, color, bold) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = txt
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.bold = bold
        p.font.name = "Calibri"
        p.space_after = Pt(8 if size < 24 else 14)
    return box


def add_brand_bar(slide, dark: bool = True) -> None:
    muted_gold = RGBColor(170, 145, 70)
    add_text(slide, 0.55, 0.35, 4, 0.3, [("ЧОО «ГЮРЗА»", 10, muted_gold, True)])
    add_text(slide, 9.8, 0.35, 3, 0.3, [("gurza36.ru", 10, MUTED, False)], align=PP_ALIGN.RIGHT)


def add_stat_row(slide, stats: list[tuple[str, str]], top: float = 2.0) -> None:
    w = 3.8
    border = RGBColor(55, 60, 72)
    for i, (val, label) in enumerate(stats):
        left = 0.7 + i * (w + 0.15)
        shape = slide.shapes.add_shape(1, Inches(left), Inches(top), Inches(w), Inches(1.55))
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(14, 18, 26)
        shape.line.color.rgb = border
        shape.line.width = Pt(0.75)
        add_text(slide, left + 0.2, top + 0.15, w - 0.4, 0.7, [(val, 36, GOLD, True)])
        add_text(slide, left + 0.2, top + 0.85, w - 0.4, 0.6, [(label, 13, MUTED, False)])


def add_bullets(slide, items: list[str], left=0.8, top=2.0, width=11.5, size=20, color=TEXT):
    box = textbox(slide, left, top, width, 4.5)
    tf = box.text_frame
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"—  {item}"
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.name = "Calibri"
        p.space_before = Pt(10)
        p.space_after = Pt(6)
        p.level = 0


def pick_bg(i: int) -> Path:
    keys = ["dark", "alt", "soft"]
    return bgs[keys[i % len(keys)]]


def add_screenshot_slide(prs, title: str, caption: str, img_path: Path, transition="fade") -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    fill_bg(slide, bgs["dark"])
    add_brand_bar(slide)
    add_text(slide, 0.7, 0.95, 10, 0.6, [(title, 32, GOLD, True)])
    add_text(slide, 0.7, 1.45, 10, 0.4, [(caption, 16, MUTED, False)])
    if img_path.exists():
        slide.shapes.add_picture(str(img_path), Inches(0.45), Inches(1.85), width=Inches(12.4))
        # frame
        frame = slide.shapes.add_shape(1, Inches(0.45), Inches(1.85), Inches(12.4), Inches(5.15))
        frame.fill.background()
        frame.line.color.rgb = GOLD
        frame.line.width = Pt(1.25)
    set_transition(slide, transition)


# --- Build ---
bgs = make_backgrounds()
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
blank = prs.slide_layouts[6]

# 01 TITLE
s = prs.slides.add_slide(blank)
fill_bg(s, bgs["dark"])
if (SCREENS / "04-logo.png").exists():
    s.shapes.add_picture(str(SCREENS / "04-logo.png"), Inches(9.2), Inches(1.5), width=Inches(3.3))
add_text(s, 0.75, 1.8, 8.5, 1.2, [("ГЮРЗА", 80, GOLD, True)])
add_text(s, 0.75, 3.0, 8.5, 1.0, [("Частная охранная", 44, TEXT, True)])
add_text(s, 0.75, 3.75, 8.5, 1.0, [("организация", 44, TEXT, False)])
add_text(s, 0.78, 5.0, 7, 0.5, [("Воронеж и Воронежская область · с 2011 года", 18, MUTED, False)])
add_text(s, 0.78, 5.55, 7, 0.5, [("Безопасность, проверенная временем", 20, GOLD_L, True)])
set_transition(s, "fade")

# 02 ABOUT
s = prs.slides.add_slide(blank)
fill_bg(s, pick_bg(1))
add_brand_bar(s)
add_text(s, 0.7, 1.0, 10, 0.5, [("О КОМПАНИИ", 14, GOLD, True)])
add_text(s, 0.7, 1.45, 11, 0.8, [("Работаем с 2011 года — и знаем цену доверию", 34, TEXT, True)])
add_text(s, 0.7, 2.5, 11, 2.2, [
    ("ООО ЧОО «Гюрза» оказывает услуги частной охраны высокого качества", 18, TEXT, False),
    ("по доступным ценам — для физических и юридических лиц.", 18, TEXT, False),
    ("", 8, TEXT, False),
    ("Разовые задачи и долгосрочные договоры. Воронеж и область.", 17, MUTED, False),
    ("", 8, TEXT, False),
    ("Гюрза — символ молниеносной реакции и абсолютного контроля территории.", 17, GOLD_L, True),
])
add_stat_row(s, [("2011", "год основания"), ("24/7", "охрана объектов"), ("100%", "опытные специалисты")], top=5.0)
set_transition(s, "fade")

# 03 SERVICES
s = prs.slides.add_slide(blank)
fill_bg(s, pick_bg(2))
add_brand_bar(s)
add_text(s, 0.7, 1.0, 10, 0.6, [("УСЛУГИ ПОД КЛЮЧ", 14, GOLD, True)])
services = [
    ("Административные здания", "Офисы, бизнес-центры, порядок и безопасность персонала"),
    ("Строительные объекты", "Площадки, техника, материалы на всех этапах"),
    ("Складские помещения", "ТМЦ, логистика, сохранность грузов"),
    ("Пропускной режим", "КПП, контроль доступа, внутриобъектовый режим"),
    ("Сельхозпредприятия", "Зерновые и масличные базы — наша специализация"),
    ("Анализ угроз", "Оценка объекта и индивидуальная система охраны"),
]
for i, (title, desc) in enumerate(services):
    col = i % 2
    row = i // 2
    left = 0.7 + col * 6.3
    top = 1.7 + row * 1.65
    shape = s.shapes.add_shape(1, Inches(left), Inches(top), Inches(6.0), Inches(1.45))
    shape.fill.solid()
    shape.fill.fore_color.rgb = PANEL
    shape.line.color.rgb = RGBColor(60, 70, 90)
    add_text(s, left + 0.15, top + 0.12, 5.6, 0.45, [(title, 17, GOLD, True)])
    add_text(s, left + 0.15, top + 0.58, 5.6, 0.7, [(desc, 13, MUTED, False)])
set_transition(s, "fade")

# 04 ADVANTAGES
s = prs.slides.add_slide(blank)
fill_bg(s, pick_bg(3))
add_brand_bar(s)
add_text(s, 0.7, 1.0, 10, 0.5, [("ПРЕИМУЩЕСТВА", 14, GOLD, True)])
add_text(s, 0.7, 1.5, 11, 0.7, [("Сотрудничая с нами, вы получаете", 32, TEXT, True)])
add_stat_row(s, [
    (f"{2026 - 2011}+", "лет на рынке охранных услуг"),
    ("100%", "специалисты высокого класса"),
    ("24/7", "непрерывная охрана"),
], top=2.5)
add_bullets(s, [
    "Любой бюджет — от экономного до высокобюджетного решения",
    "Индивидуальный подход к каждому объекту",
    "Оперативная реакция и чёткий контроль территории",
    "Профессиональная оценка рисков перед началом работы",
], top=4.3, size=18)
set_transition(s, "fade")

# 05 SPECIALIZATION
s = prs.slides.add_slide(blank)
fill_bg(s, pick_bg(4))
add_brand_bar(s)
add_text(s, 0.7, 1.0, 10, 0.5, [("СПЕЦИАЛИЗАЦИЯ", 14, GOLD, True)])
add_text(s, 0.7, 1.45, 11, 0.8, [("Охрана агропредприятий и зерновых баз", 34, TEXT, True)])
add_text(s, 0.7, 2.5, 6.5, 2.5, [
    ("Особая компетенция — охрана сельхозпредприятий и баз", 18, TEXT, False),
    ("по приёмке, подработке, хранению и отпуску", 18, TEXT, False),
    ("зерновых и масличной продукции.", 18, GOLD_L, True),
    ("", 10, TEXT, False),
    ("Понимаем специфику сезонных рисков, логистики и", 17, MUTED, False),
    ("сохранности урожая — работаем точечно под задачу.", 17, MUTED, False),
])
if (SCREENS / "02-about.png").exists():
    s.shapes.add_picture(str(SCREENS / "02-about.png"), Inches(7.3), Inches(1.8), width=Inches(5.5))
set_transition(s, "fade")

# 06 GEO
s = prs.slides.add_slide(blank)
fill_bg(s, pick_bg(5))
add_brand_bar(s)
add_text(s, 0.7, 1.0, 10, 0.5, [("ГЕОГРАФИЯ", 14, GOLD, True)])
add_text(s, 0.7, 1.45, 11, 0.7, [("Воронеж и Воронежская область", 36, TEXT, True)])
add_bullets(s, [
    "Административные и офисные здания",
    "Строительные площадки и объекты на всех стадиях",
    "Склады, логистические и производственные комплексы",
    "Сельхозпредприятия, элеваторы, зерновые базы",
    "Объекты с пропускным и внутриобъектовым режимом",
], top=2.4, size=19)
set_transition(s, "fade")

# 07 CONTACTS
s = prs.slides.add_slide(blank)
fill_bg(s, pick_bg(6))
add_brand_bar(s)
add_text(s, 0.7, 1.0, 10, 0.5, [("КОНТАКТЫ", 14, GOLD, True)])
add_text(s, 0.7, 1.45, 11, 0.7, [("Связаться с нами", 36, TEXT, True)])
add_text(s, 0.7, 2.3, 5.8, 2.8, [
    ("Официальный сайт:", 16, MUTED, False),
    ("gurza36.ru", 28, GOLD_L, True),
    ("", 8, TEXT, False),
    ("Телефоны:", 16, MUTED, False),
    ("+7 (473) 278-51-47", 20, TEXT, True),
    ("+7 (900) 954-43-12", 20, TEXT, True),
    ("", 8, TEXT, False),
    ("Email: mike1592@yandex.ru", 17, TEXT, False),
])
add_text(s, 0.7, 5.2, 6, 1.2, [
    ("394016, г. Воронеж,", 16, MUTED, False),
    ("ул. Солнечная, 17/1, офис 4", 16, MUTED, False),
])
if (SCREENS / "03-contacts.png").exists():
    s.shapes.add_picture(str(SCREENS / "03-contacts.png"), Inches(6.8), Inches(1.9), width=Inches(6.0))
set_transition(s, "fade")

# 08 DIRECTOR
s = prs.slides.add_slide(blank)
fill_bg(s, pick_bg(7))
add_brand_bar(s)
add_text(s, 0.7, 1.0, 10, 0.5, [("РУКОВОДСТВО", 14, GOLD, True)])
add_text(s, 0.7, 1.8, 11, 0.8, [("Самофалов Михаил Иванович", 36, TEXT, True)])
add_text(s, 0.7, 2.7, 11, 0.5, [("Директор ООО ЧОО «Гюрза»", 22, GOLD, True)])
add_text(s, 0.7, 3.5, 10, 1.5, [
    ("Готовы обсудить охрану вашего объекта,", 18, TEXT, False),
    ("провести экспертную оценку и предложить", 18, TEXT, False),
    ("оптимальный вариант охраны под ваш бюджет.", 18, TEXT, False),
])
add_bullets(s, [
    "Звоните — проконсультируем и ответим на все вопросы",
    "Выезд на объект для оценки угроз",
    "Гибкие условия: разовые и долгосрочные договоры",
], top=4.8, size=18)
set_transition(s, "fade")

# 09 FINALE
s = prs.slides.add_slide(blank)
fill_bg(s, bgs["dark"])
if (SCREENS / "04-logo.png").exists():
    s.shapes.add_picture(str(SCREENS / "04-logo.png"), Inches(5.3), Inches(0.8), width=Inches(2.6))
add_text(s, 0.75, 3.2, 12, 0.8, [("ЧОО «Гюрза»", 48, GOLD, True)], align=PP_ALIGN.CENTER)
add_text(s, 0.75, 4.1, 12, 1.5, [
    ("+7 (473) 278-51-47  ·  +7 (900) 954-43-12", 18, TEXT, False),
    ("mike1592@yandex.ru  ·  gurza36.ru", 18, GOLD_L, False),
    ("394016, г. Воронеж, ул. Солнечная, 17/1, офис 4", 16, MUTED, False),
], align=PP_ALIGN.CENTER)
set_transition(s, "fade")

prs.save(OUT)
print(f"Wrote {OUT} ({len(prs.slides)} slides)")
